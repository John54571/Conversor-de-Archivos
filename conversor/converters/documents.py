import os
from pathlib import Path

from .base import BaseConverter, ConversionTask, ConversionOptions, FileCategory
from ..core.registry import ConverterRegistry
from ..utils.config import config_manager
from ..utils.logger import logger


@ConverterRegistry.register
class DocumentConverter(BaseConverter):
    SUPPORTED_INPUT = [
        "docx", "pdf", "xlsx", "csv", "pptx", "txt", "rtf", "html", "epub", "odt"
    ]
    SUPPORTED_OUTPUT = [
        "docx", "pdf", "xlsx", "csv", "txt", "html", "epub", "odt"
    ]
    CATEGORY = FileCategory.DOCUMENT

    BLOCKED_CONVERSIONS = {
        ("xlsx", "docx"), ("csv", "docx"),
        ("docx", "xlsx"), ("docx", "csv"),
        ("pdf", "xlsx"), ("pdf", "csv"),
        ("pptx", "xlsx"), ("pptx", "csv"),
    }

    _UNICODE_REPLACEMENTS = {
        "\u201c": '"', "\u201d": '"',
        "\u2018": "'", "\u2019": "'",
        "\u2013": "-", "\u2014": "--",
        "\u2026": "...",
        "\u00a0": " ",
        "\u00ab": '"', "\u00bb": '"',
        "\u201a": "'", "\u201b": "'",
        "\u201e": '"', "\u201f": '"',
    }

    @staticmethod
    def _sanitize_text(text: str) -> str:
        for char, replacement in DocumentConverter._UNICODE_REPLACEMENTS.items():
            text = text.replace(char, replacement)
        text = text.encode("latin-1", errors="replace").decode("latin-1")
        return text

    @classmethod
    def can_convert(cls, source_ext: str, target_ext: str) -> bool:
        src = source_ext.lower().lstrip(".")
        tgt = target_ext.lower().lstrip(".")
        if (src, tgt) in cls.BLOCKED_CONVERSIONS:
            return False
        if src == tgt:
            return False
        return (src in cls.SUPPORTED_INPUT and tgt in cls.SUPPORTED_OUTPUT)

    @classmethod
    def get_valid_outputs(cls, source_ext: str) -> list[str]:
        ext = source_ext.lower().lstrip(".")
        if ext not in cls.SUPPORTED_INPUT:
            return []
        outputs = []
        for fmt in cls.SUPPORTED_OUTPUT:
            if fmt != ext and (ext, fmt) not in cls.BLOCKED_CONVERSIONS:
                outputs.append(fmt)
        return outputs

    def convert(self, task: ConversionTask) -> Path:
        src = task.source_path
        src_ext = src.suffix.lower().lstrip(".")
        tgt_ext = task.output_format.lower()
        output = task.options.output_dir or str(src.parent)
        output_path = Path(output) / f"{src.stem}.{tgt_ext}"

        Path(output).mkdir(parents=True, exist_ok=True)

        converters = {
            ("docx", "pdf"): self._docx_to_pdf,
            ("pdf", "docx"): self._pdf_to_docx,
            ("xlsx", "csv"): self._xlsx_to_csv,
            ("csv", "xlsx"): self._csv_to_xlsx,
            ("txt", "pdf"): self._txt_to_pdf,
            ("html", "pdf"): self._html_to_pdf,
            ("docx", "txt"): self._docx_to_txt,
            ("pdf", "txt"): self._pdf_to_txt,
            ("txt", "docx"): self._txt_to_docx,
            ("txt", "html"): self._txt_to_html,
            ("html", "txt"): self._html_to_txt,
            ("docx", "html"): self._docx_to_html,
            ("html", "docx"): self._html_to_docx,
            ("pptx", "pdf"): self._pptx_to_pdf,
            ("csv", "txt"): self._csv_to_txt,
            ("xlsx", "txt"): self._xlsx_to_txt,
            ("rtf", "txt"): self._rtf_to_txt,
            ("rtf", "pdf"): self._rtf_to_pdf,
            ("rtf", "docx"): self._rtf_to_docx,
            ("rtf", "html"): self._rtf_to_html,
            ("epub", "txt"): self._epub_to_txt,
            ("epub", "html"): self._epub_to_html,
            ("epub", "pdf"): self._epub_to_pdf,
            ("odt", "txt"): self._odt_to_txt,
            ("odt", "pdf"): self._odt_to_pdf,
            ("odt", "docx"): self._odt_to_docx,
            ("docx", "odt"): self._docx_to_odt,
            ("txt", "csv"): self._txt_to_csv,
            ("txt", "xlsx"): self._txt_to_xlsx,
            ("txt", "odt"): self._txt_to_odt,
            ("txt", "epub"): self._txt_to_epub,
            ("csv", "html"): self._csv_to_html,
            ("csv", "pdf"): self._csv_to_pdf,
            ("csv", "odt"): self._csv_to_odt,
            ("csv", "epub"): self._csv_to_epub,
            ("xlsx", "html"): self._xlsx_to_html,
            ("xlsx", "pdf"): self._xlsx_to_pdf,
            ("xlsx", "odt"): self._xlsx_to_odt,
            ("xlsx", "epub"): self._xlsx_to_epub,
            ("html", "csv"): self._html_to_csv,
            ("html", "xlsx"): self._html_to_xlsx,
            ("html", "odt"): self._html_to_odt,
            ("html", "epub"): self._html_to_epub,
            ("rtf", "csv"): self._rtf_to_csv,
            ("rtf", "xlsx"): self._rtf_to_xlsx,
            ("rtf", "odt"): self._rtf_to_odt,
            ("rtf", "epub"): self._rtf_to_epub,
            ("odt", "csv"): self._odt_to_csv,
            ("odt", "xlsx"): self._odt_to_xlsx,
            ("odt", "html"): self._odt_to_html,
            ("odt", "epub"): self._odt_to_epub,
            ("epub", "csv"): self._epub_to_csv,
            ("epub", "xlsx"): self._epub_to_xlsx,
            ("epub", "odt"): self._epub_to_odt,
            ("epub", "docx"): self._epub_to_docx,
            ("pdf", "html"): self._pdf_to_html,
            ("pdf", "odt"): self._pdf_to_odt,
            ("pdf", "epub"): self._pdf_to_epub,
            ("pdf", "xlsx"): self._pdf_to_xlsx,
            ("pptx", "txt"): self._pptx_to_txt,
            ("pptx", "docx"): self._pptx_to_docx,
            ("pptx", "html"): self._pptx_to_html,
            ("pptx", "odt"): self._pptx_to_odt,
            ("pptx", "epub"): self._pptx_to_epub,
            ("docx", "epub"): self._docx_to_epub,
        }

        handler = converters.get((src_ext, tgt_ext))
        if handler:
            handler(src, output_path, task)
        else:
            raise ValueError(f"Conversion no soportada: {src_ext} -> {tgt_ext}")

        return output_path

    def _docx_to_pdf(self, src: Path, dst: Path, task: ConversionTask):
        from docx import Document
        from docx.oxml.text.paragraph import CT_P
        from docx.oxml.table import CT_Tbl
        from fpdf import FPDF

        doc = Document(str(src))
        pdf = FPDF()
        pdf.set_auto_page_break(auto=True, margin=15)
        pdf.add_page()
        pdf.set_font("Helvetica", size=11)

        for element in doc.element.body:
            if isinstance(element, CT_P):
                para = None
                for p in doc.paragraphs:
                    if p._element is element:
                        para = p
                        break
                if para:
                    text = self._sanitize_text(para.text)
                    if text.strip():
                        pdf.multi_cell(w=190, h=6, text=text)
                    else:
                        pdf.ln(4)
            elif isinstance(element, CT_Tbl):
                table = None
                for t in doc.tables:
                    if t._tbl is element:
                        table = t
                        break
                if table:
                    pdf.set_font("Helvetica", size=9)
                    for row in table.rows:
                        row_text = " | ".join(self._sanitize_text(cell.text) for cell in row.cells)
                        if row_text.strip():
                            pdf.multi_cell(w=190, h=5, text=row_text)
                    pdf.set_font("Helvetica", size=11)
                    pdf.ln(2)

        if task.on_progress:
            task.on_progress(0.5)

        pdf.output(str(dst))

        if task.on_progress:
            task.on_progress(1.0)

    def _pdf_to_docx(self, src: Path, dst: Path, task: ConversionTask):
        from pdf2docx import Converter as Pdf2DocxConverter

        cv = Pdf2DocxConverter(str(src))
        cv.convert(str(dst))
        cv.close()

        if task.on_progress:
            task.on_progress(1.0)

    def _xlsx_to_csv(self, src: Path, dst: Path, task: ConversionTask):
        import pandas as pd

        df = pd.read_excel(str(src), sheet_name=0)
        df.to_csv(str(dst), index=False, encoding="utf-8")

        if task.on_progress:
            task.on_progress(1.0)

    def _csv_to_xlsx(self, src: Path, dst: Path, task: ConversionTask):
        import pandas as pd

        df = pd.read_csv(str(src))
        df.to_excel(str(dst), index=False, engine="openpyxl")

        if task.on_progress:
            task.on_progress(1.0)

    def _txt_to_pdf(self, src: Path, dst: Path, task: ConversionTask):
        from fpdf import FPDF

        text = src.read_text(encoding="utf-8")
        pdf = FPDF()
        pdf.set_auto_page_break(auto=True, margin=15)
        pdf.add_page()
        pdf.set_font("Courier", size=10)

        for line in text.split("\n"):
            pdf.multi_cell(w=190, h=5, text=self._sanitize_text(line))

        if task.on_progress:
            task.on_progress(0.7)

        pdf.output(str(dst))

        if task.on_progress:
            task.on_progress(1.0)

    def _html_to_pdf(self, src: Path, dst: Path, task: ConversionTask):
        from fpdf import FPDF
        import re

        html_content = src.read_text(encoding="utf-8")
        text = re.sub(r"<[^>]+>", "", html_content)
        text = re.sub(r"\s+", " ", text).strip()

        pdf = FPDF()
        pdf.set_auto_page_break(auto=True, margin=15)
        pdf.add_page()
        pdf.set_font("Helvetica", size=11)
        pdf.multi_cell(w=190, h=6, text=self._sanitize_text(text))

        if task.on_progress:
            task.on_progress(0.7)

        pdf.output(str(dst))

        if task.on_progress:
            task.on_progress(1.0)

    def _docx_to_txt(self, src: Path, dst: Path, task: ConversionTask):
        from docx import Document

        doc = Document(str(src))
        text = "\n".join(para.text for para in doc.paragraphs)
        dst.write_text(text, encoding="utf-8")

        if task.on_progress:
            task.on_progress(1.0)

    def _pdf_to_txt(self, src: Path, dst: Path, task: ConversionTask):
        use_ocr = config_manager.get_config().enable_ocr

        if use_ocr:
            self._pdf_to_txt_ocr(src, dst, task)
        else:
            self._pdf_to_txt_basic(src, dst, task)

    def _pdf_to_txt_basic(self, src: Path, dst: Path, task: ConversionTask):
        from pypdf import PdfReader

        reader = PdfReader(str(src))
        text_parts = []
        total = len(reader.pages)

        for i, page in enumerate(reader.pages):
            text_parts.append(page.extract_text() or "")
            if task.on_progress:
                task.on_progress((i + 1) / total * 0.9)

        dst.write_text("\n\n".join(text_parts), encoding="utf-8")

        if task.on_progress:
            task.on_progress(1.0)

    def _pdf_to_txt_ocr(self, src: Path, dst: Path, task: ConversionTask):
        try:
            from pdf2image import convert_from_path
            from paddleocr import PaddleOCR
            import gc

            ocr = PaddleOCR(use_angle_cls=True, lang=config_manager.get_config().ocr_language)
            
            from pypdf import PdfReader
            reader = PdfReader(str(src))
            total_pages = len(reader.pages)
            
            if total_pages > 50:
                ram_gb = self._get_available_ram_gb()
                if ram_gb < 4:
                    logger.warning(f"PDF grande ({total_pages} paginas) con RAM limitada ({ram_gb:.1f}GB). Puede ser lento.")
            
            text_parts = []
            for i in range(total_pages):
                images = convert_from_path(str(src), first_page=i+1, last_page=i+1, dpi=150)
                if not images:
                    text_parts.append("")
                    continue
                
                img = images[0]
                img_path = src.parent / f"_temp_page_{i}.png"
                img.save(str(img_path))
                del images
                del img
                
                result = ocr.ocr(str(img_path), cls=True)
                page_text = ""
                if result and result[0]:
                    for line in result[0]:
                        if line and len(line) > 1:
                            page_text += line[1][0] + "\n"
                
                text_parts.append(page_text)
                
                if img_path.exists():
                    img_path.unlink()
                
                gc.collect()
                
                if task.on_progress:
                    task.on_progress((i + 1) / total_pages * 0.9)

            dst.write_text("\n\n".join(text_parts), encoding="utf-8")

            if task.on_progress:
                task.on_progress(1.0)

        except ImportError as e:
            logger.warning(f"PaddleOCR no disponible, usando extracción básica: {e}")
            self._pdf_to_txt_basic(src, dst, task)
        except Exception as e:
            logger.error(f"Error en OCR: {e}")
            raise ValueError(f"Error en OCR: {e}")

    def _get_available_ram_gb(self) -> float:
        try:
            import ctypes
            import sys
            if sys.platform == "win32":
                class MEMORYSTATUSEX(ctypes.Structure):
                    _fields_ = [
                        ("dwLength", ctypes.c_ulong),
                        ("dwMemoryLoad", ctypes.c_ulong),
                        ("ullTotalPhys", ctypes.c_ulonglong),
                        ("ullAvailPhys", ctypes.c_ulonglong),
                        ("ullTotalPageFile", ctypes.c_ulonglong),
                        ("ullAvailPageFile", ctypes.c_ulonglong),
                        ("ullTotalVirtual", ctypes.c_ulonglong),
                        ("ullAvailVirtual", ctypes.c_ulonglong),
                        ("sullAvailExtendedVirtual", ctypes.c_ulonglong),
                    ]
                    def __init__(self):
                        self.dwLength = ctypes.sizeof(self)
                
                stat = MEMORYSTATUSEX()
                ctypes.windll.kernel32.GlobalMemoryStatusEx(ctypes.byref(stat))
                return stat.ullAvailPhys / (1024 ** 3)
        except Exception:
            pass
        return 8.0

    def _txt_to_docx(self, src: Path, dst: Path, task: ConversionTask):
        from docx import Document

        text = src.read_text(encoding="utf-8")
        doc = Document()
        for line in text.split("\n"):
            doc.add_paragraph(line)
        doc.save(str(dst))

        if task.on_progress:
            task.on_progress(1.0)

    def _txt_to_html(self, src: Path, dst: Path, task: ConversionTask):
        text = src.read_text(encoding="utf-8")
        html = f"<!DOCTYPE html>\n<html>\n<head><meta charset='utf-8'><title>{src.stem}</title></head>\n<body>\n"
        for line in text.split("\n"):
            html += f"<p>{line}</p>\n" if line.strip() else "<br>\n"
        html += "</body>\n</html>"
        dst.write_text(html, encoding="utf-8")

        if task.on_progress:
            task.on_progress(1.0)

    def _html_to_txt(self, src: Path, dst: Path, task: ConversionTask):
        import re

        html = src.read_text(encoding="utf-8")
        text = re.sub(r"<br\s*/?>", "\n", html)
        text = re.sub(r"</p>", "\n", text)
        text = re.sub(r"<[^>]+>", "", text)
        text = re.sub(r"\n{3,}", "\n\n", text)
        dst.write_text(text.strip(), encoding="utf-8")

        if task.on_progress:
            task.on_progress(1.0)

    def _docx_to_html(self, src: Path, dst: Path, task: ConversionTask):
        from docx import Document

        doc = Document(str(src))
        html = f"<!DOCTYPE html>\n<html>\n<head><meta charset='utf-8'><title>{src.stem}</title></head>\n<body>\n"
        for para in doc.paragraphs:
            if para.text.strip():
                html += f"<p>{para.text}</p>\n"
        html += "</body>\n</html>"
        dst.write_text(html, encoding="utf-8")

        if task.on_progress:
            task.on_progress(1.0)

    def _html_to_docx(self, src: Path, dst: Path, task: ConversionTask):
        import re
        from docx import Document

        html = src.read_text(encoding="utf-8")
        text = re.sub(r"<br\s*/?>", "\n", html)
        text = re.sub(r"</p>", "\n", text)
        text = re.sub(r"<[^>]+>", "", text)

        doc = Document()
        for line in text.split("\n"):
            if line.strip():
                doc.add_paragraph(line.strip())
        doc.save(str(dst))

        if task.on_progress:
            task.on_progress(1.0)

    def _pptx_to_pdf(self, src: Path, dst: Path, task: ConversionTask):
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from fpdf import FPDF
        from PIL import Image
        import io

        prs = Presentation(str(src))
        pdf = FPDF()
        total = len(prs.slides)

        for i, slide in enumerate(prs.slides):
            pdf.add_page()
            
            pdf.set_font("Helvetica", "B", 16)
            pdf.cell(0, 10, f"Slide {i + 1}", ln=True)
            pdf.ln(5)

            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        if para.text.strip():
                            pdf.set_font("Helvetica", size=12)
                            pdf.multi_cell(w=190, h=6, text=self._sanitize_text(para.text))
                            pdf.ln(2)
                
                if shape.shape_type == 13:
                    try:
                        image = shape.image
                        image_bytes = image.blob
                        img = Image.open(io.BytesIO(image_bytes))
                        
                        temp_img_path = src.parent / f"_temp_slide_img_{i}_{id(shape)}.png"
                        img.save(str(temp_img_path))
                        
                        pdf.image(str(temp_img_path), x=10, y=pdf.get_y(), w=100)
                        pdf.ln(img.height * 0.2 + 5)
                        
                        if temp_img_path.exists():
                            temp_img_path.unlink()
                    except Exception as e:
                        logger.debug(f"No se pudo insertar imagen en slide {i+1}: {e}")

            if task.on_progress:
                task.on_progress((i + 1) / total * 0.9)

        pdf.output(str(dst))

        if task.on_progress:
            task.on_progress(1.0)

    def _csv_to_txt(self, src: Path, dst: Path, task: ConversionTask):
        text = src.read_text(encoding="utf-8")
        dst.write_text(text, encoding="utf-8")

        if task.on_progress:
            task.on_progress(1.0)

    def _xlsx_to_txt(self, src: Path, dst: Path, task: ConversionTask):
        import pandas as pd

        df = pd.read_excel(str(src), sheet_name=0)
        dst.write_text(df.to_string(index=False), encoding="utf-8")

        if task.on_progress:
            task.on_progress(1.0)

    def _rtf_to_txt(self, src: Path, dst: Path, task: ConversionTask):
        import re

        rtf = src.read_text(encoding="utf-8", errors="ignore")
        text = re.sub(r"\\par\b", "\n", rtf)
        text = re.sub(r"\{[^}]*\}", "", text)
        text = re.sub(r"\\\w+\d*\s?", "", text)
        text = re.sub(r"[{}]", "", text)
        text = re.sub(r"\n{3,}", "\n\n", text)
        dst.write_text(text.strip(), encoding="utf-8")

        if task.on_progress:
            task.on_progress(1.0)

    def _rtf_to_pdf(self, src: Path, dst: Path, task: ConversionTask):
        import re
        from fpdf import FPDF

        rtf = src.read_text(encoding="utf-8", errors="ignore")
        text = re.sub(r"\\par\b", "\n", rtf)
        text = re.sub(r"\{[^}]*\}", "", text)
        text = re.sub(r"\\\w+\d*\s?", "", text)
        text = re.sub(r"[{}]", "", text)
        text = text.strip()

        pdf = FPDF()
        pdf.set_auto_page_break(auto=True, margin=15)
        pdf.add_page()
        pdf.set_font("Helvetica", size=11)
        pdf.multi_cell(w=190, h=6, text=self._sanitize_text(text))

        if task.on_progress:
            task.on_progress(0.7)

        pdf.output(str(dst))

        if task.on_progress:
            task.on_progress(1.0)

    def _rtf_to_docx(self, src: Path, dst: Path, task: ConversionTask):
        import re
        from docx import Document

        rtf = src.read_text(encoding="utf-8", errors="ignore")
        text = re.sub(r"\\par\b", "\n", rtf)
        text = re.sub(r"\{[^}]*\}", "", text)
        text = re.sub(r"\\\w+\d*\s?", "", text)
        text = re.sub(r"[{}]", "", text)
        text = text.strip()

        doc = Document()
        for line in text.split("\n"):
            if line.strip():
                doc.add_paragraph(line.strip())
        doc.save(str(dst))

        if task.on_progress:
            task.on_progress(1.0)

    def _rtf_to_html(self, src: Path, dst: Path, task: ConversionTask):
        import re

        rtf = src.read_text(encoding="utf-8", errors="ignore")
        text = re.sub(r"\\par\b", "\n", rtf)
        text = re.sub(r"\{[^}]*\}", "", text)
        text = re.sub(r"\\\w+\d*\s?", "", text)
        text = re.sub(r"[{}]", "", text)
        text = text.strip()

        html = f"<!DOCTYPE html>\n<html>\n<head><meta charset='utf-8'><title>{src.stem}</title></head>\n<body>\n"
        for line in text.split("\n"):
            html += f"<p>{line}</p>\n" if line.strip() else "<br>\n"
        html += "</body>\n</html>"
        dst.write_text(html, encoding="utf-8")

        if task.on_progress:
            task.on_progress(1.0)

    def _epub_to_txt(self, src: Path, dst: Path, task: ConversionTask):
        from ebooklib import epub
        from bs4 import BeautifulSoup

        book = epub.read_epub(str(src))
        text_parts = []
        items = list(book.get_items_of_type(9))
        total = len(items)

        for i, item in enumerate(items):
            html_content = item.get_content().decode("utf-8")
            soup = BeautifulSoup(html_content, "html.parser")
            text = soup.get_text()
            text_parts.append(text)
            
            if task.on_progress:
                task.on_progress((i + 1) / total * 0.9)

        dst.write_text("\n\n".join(text_parts), encoding="utf-8")

        if task.on_progress:
            task.on_progress(1.0)

    def _epub_to_html(self, src: Path, dst: Path, task: ConversionTask):
        from ebooklib import epub

        book = epub.read_epub(str(src))
        html_parts = []
        items = list(book.get_items_of_type(9))
        total = len(items)

        for i, item in enumerate(items):
            html_content = item.get_content().decode("utf-8")
            html_parts.append(html_content)
            
            if task.on_progress:
                task.on_progress((i + 1) / total * 0.9)

        combined = "\n\n".join(html_parts)
        dst.write_text(combined, encoding="utf-8")

        if task.on_progress:
            task.on_progress(1.0)

    def _epub_to_pdf(self, src: Path, dst: Path, task: ConversionTask):
        from ebooklib import epub
        from bs4 import BeautifulSoup
        from fpdf import FPDF

        book = epub.read_epub(str(src))
        pdf = FPDF()
        pdf.set_auto_page_break(auto=True, margin=15)
        pdf.add_page()
        pdf.set_font("Helvetica", size=11)

        items = list(book.get_items_of_type(9))
        total = len(items)

        for i, item in enumerate(items):
            html_content = item.get_content().decode("utf-8")
            soup = BeautifulSoup(html_content, "html.parser")
            text = soup.get_text()
            
            for line in text.split("\n"):
                if line.strip():
                    pdf.multi_cell(w=190, h=6, text=self._sanitize_text(line.strip()))
            
            if task.on_progress:
                task.on_progress((i + 1) / total * 0.9)

        pdf.output(str(dst))

        if task.on_progress:
            task.on_progress(1.0)

    def _odt_to_txt(self, src: Path, dst: Path, task: ConversionTask):
        import zipfile
        from xml.etree import ElementTree as ET

        with zipfile.ZipFile(str(src), "r") as z:
            content_xml = z.read("content.xml")
        
        root = ET.fromstring(content_xml)
        text_parts = []
        
        for elem in root.iter():
            if elem.text and elem.text.strip():
                text_parts.append(elem.text.strip())

        dst.write_text("\n".join(text_parts), encoding="utf-8")

        if task.on_progress:
            task.on_progress(1.0)

    def _odt_to_pdf(self, src: Path, dst: Path, task: ConversionTask):
        import zipfile
        from xml.etree import ElementTree as ET
        from fpdf import FPDF

        with zipfile.ZipFile(str(src), "r") as z:
            content_xml = z.read("content.xml")
        
        root = ET.fromstring(content_xml)
        text_parts = []
        
        for elem in root.iter():
            if elem.text and elem.text.strip():
                text_parts.append(elem.text.strip())

        pdf = FPDF()
        pdf.set_auto_page_break(auto=True, margin=15)
        pdf.add_page()
        pdf.set_font("Helvetica", size=11)

        for text in text_parts:
            pdf.multi_cell(w=190, h=6, text=self._sanitize_text(text))

        if task.on_progress:
            task.on_progress(0.7)

        pdf.output(str(dst))

        if task.on_progress:
            task.on_progress(1.0)

    def _odt_to_docx(self, src: Path, dst: Path, task: ConversionTask):
        import zipfile
        from xml.etree import ElementTree as ET
        from docx import Document

        with zipfile.ZipFile(str(src), "r") as z:
            content_xml = z.read("content.xml")
        
        root = ET.fromstring(content_xml)
        text_parts = []
        
        for elem in root.iter():
            if elem.text and elem.text.strip():
                text_parts.append(elem.text.strip())

        doc = Document()
        for text in text_parts:
            doc.add_paragraph(text)
        doc.save(str(dst))

        if task.on_progress:
            task.on_progress(1.0)

    def _docx_to_odt(self, src: Path, dst: Path, task: ConversionTask):
        from docx import Document
        from odf.opendocument import OpenDocumentText
        from odf.text import P
        from odf.table import Table, TableRow, TableCell

        doc = Document(str(src))
        odt = OpenDocumentText()

        for para in doc.paragraphs:
            if para.text.strip():
                odt.text.addElement(P(text=para.text))
            else:
                odt.text.addElement(P(text=""))

        for table in doc.tables:
            odf_table = Table()
            for row in table.rows:
                odf_row = TableRow()
                for cell in row.cells:
                    odf_cell = TableCell()
                    odf_cell.addElement(P(text=cell.text))
                    odf_row.addElement(odf_cell)
                odf_table.addElement(odf_row)
            odt.text.addElement(odf_table)

        odt.save(str(dst))

        if task.on_progress:
            task.on_progress(1.0)

    def _txt_to_csv(self, src: Path, dst: Path, task: ConversionTask):
        import csv
        text = src.read_text(encoding="utf-8")
        with open(dst, "w", encoding="utf-8", newline="") as f:
            writer = csv.writer(f)
            for line in text.split("\n"):
                if line.strip():
                    writer.writerow([line])
        if task.on_progress:
            task.on_progress(1.0)

    def _txt_to_xlsx(self, src: Path, dst: Path, task: ConversionTask):
        import pandas as pd
        text = src.read_text(encoding="utf-8")
        lines = [line for line in text.split("\n") if line.strip()]
        df = pd.DataFrame({"Contenido": lines})
        df.to_excel(str(dst), index=False, engine="openpyxl")
        if task.on_progress:
            task.on_progress(1.0)

    def _txt_to_odt(self, src: Path, dst: Path, task: ConversionTask):
        from odf.opendocument import OpenDocumentText
        from odf.text import P
        text = src.read_text(encoding="utf-8")
        odt = OpenDocumentText()
        for line in text.split("\n"):
            odt.text.addElement(P(text=line))
        odt.save(str(dst))
        if task.on_progress:
            task.on_progress(1.0)

    def _txt_to_epub(self, src: Path, dst: Path, task: ConversionTask):
        from ebooklib import epub
        text = src.read_text(encoding="utf-8")
        book = epub.EpubBook()
        book.set_identifier("id-1234")
        book.set_title(src.stem)
        book.set_language("es")
        chapter = epub.EpubHtml(title="Contenido", file_name="content.xhtml", lang="es")
        html_content = "<html><body>"
        for line in text.split("\n"):
            if line.strip():
                html_content += f"<p>{line}</p>"
            else:
                html_content += "<br/>"
        html_content += "</body></html>"
        chapter.content = html_content
        book.add_item(chapter)
        book.spine = ["nav", chapter]
        epub.write_epub(str(dst), book)
        if task.on_progress:
            task.on_progress(1.0)

    def _csv_to_html(self, src: Path, dst: Path, task: ConversionTask):
        import csv
        with open(src, "r", encoding="utf-8") as f:
            reader = csv.reader(f)
            rows = list(reader)
        html = f"<!DOCTYPE html><html><head><meta charset='utf-8'><title>{src.stem}</title></head><body><table border='1'>"
        for i, row in enumerate(rows):
            html += "<tr>"
            tag = "th" if i == 0 else "td"
            for cell in row:
                html += f"<{tag}>{cell}</{tag}>"
            html += "</tr>"
        html += "</table></body></html>"
        dst.write_text(html, encoding="utf-8")
        if task.on_progress:
            task.on_progress(1.0)

    def _csv_to_pdf(self, src: Path, dst: Path, task: ConversionTask):
        import csv
        from fpdf import FPDF
        with open(src, "r", encoding="utf-8") as f:
            reader = csv.reader(f)
            rows = list(reader)
        pdf = FPDF()
        pdf.set_auto_page_break(auto=True, margin=15)
        pdf.add_page()
        pdf.set_font("Helvetica", size=9)
        for row in rows:
            line = " | ".join(self._sanitize_text(cell) for cell in row)
            if line.strip():
                pdf.multi_cell(w=190, h=5, text=line)
        pdf.output(str(dst))
        if task.on_progress:
            task.on_progress(1.0)

    def _csv_to_odt(self, src: Path, dst: Path, task: ConversionTask):
        import csv
        from odf.opendocument import OpenDocumentText
        from odf.text import P
        from odf.table import Table, TableRow, TableCell
        with open(src, "r", encoding="utf-8") as f:
            reader = csv.reader(f)
            rows = list(reader)
        odt = OpenDocumentText()
        if rows:
            table = Table()
            for row in rows:
                odf_row = TableRow()
                for cell in row:
                    odf_cell = TableCell()
                    odf_cell.addElement(P(text=cell))
                    odf_row.addElement(odf_cell)
                table.addElement(odf_row)
            odt.text.addElement(table)
        odt.save(str(dst))
        if task.on_progress:
            task.on_progress(1.0)

    def _csv_to_epub(self, src: Path, dst: Path, task: ConversionTask):
        import csv
        from ebooklib import epub
        with open(src, "r", encoding="utf-8") as f:
            reader = csv.reader(f)
            rows = list(reader)
        book = epub.EpubBook()
        book.set_identifier("id-1234")
        book.set_title(src.stem)
        book.set_language("es")
        chapter = epub.EpubHtml(title="Contenido", file_name="content.xhtml", lang="es")
        html_content = "<html><body><table border='1'>"
        for i, row in enumerate(rows):
            html_content += "<tr>"
            tag = "th" if i == 0 else "td"
            for cell in row:
                html_content += f"<{tag}>{cell}</{tag}>"
            html_content += "</tr>"
        html_content += "</table></body></html>"
        chapter.content = html_content
        book.add_item(chapter)
        book.spine = ["nav", chapter]
        epub.write_epub(str(dst), book)
        if task.on_progress:
            task.on_progress(1.0)

    def _xlsx_to_html(self, src: Path, dst: Path, task: ConversionTask):
        import pandas as pd
        df = pd.read_excel(str(src), sheet_name=0)
        html = df.to_html(index=False)
        full_html = f"<!DOCTYPE html><html><head><meta charset='utf-8'><title>{src.stem}</title></head><body>{html}</body></html>"
        dst.write_text(full_html, encoding="utf-8")
        if task.on_progress:
            task.on_progress(1.0)

    def _xlsx_to_pdf(self, src: Path, dst: Path, task: ConversionTask):
        import pandas as pd
        from fpdf import FPDF
        df = pd.read_excel(str(src), sheet_name=0)
        pdf = FPDF()
        pdf.set_auto_page_break(auto=True, margin=15)
        pdf.add_page()
        pdf.set_font("Helvetica", size=9)
        for _, row in df.iterrows():
            line = " | ".join(self._sanitize_text(str(cell)) for cell in row)
            if line.strip():
                pdf.multi_cell(w=190, h=5, text=line)
        pdf.output(str(dst))
        if task.on_progress:
            task.on_progress(1.0)

    def _xlsx_to_odt(self, src: Path, dst: Path, task: ConversionTask):
        import pandas as pd
        from odf.opendocument import OpenDocumentText
        from odf.text import P
        from odf.table import Table, TableRow, TableCell
        df = pd.read_excel(str(src), sheet_name=0)
        odt = OpenDocumentText()
        table = Table()
        header_row = TableRow()
        for col in df.columns:
            cell = TableCell()
            cell.addElement(P(text=str(col)))
            header_row.addElement(cell)
        table.addElement(header_row)
        for _, row in df.iterrows():
            odf_row = TableRow()
            for cell in row:
                odf_cell = TableCell()
                odf_cell.addElement(P(text=str(cell)))
                odf_row.addElement(odf_cell)
            table.addElement(odf_row)
        odt.text.addElement(table)
        odt.save(str(dst))
        if task.on_progress:
            task.on_progress(1.0)

    def _xlsx_to_epub(self, src: Path, dst: Path, task: ConversionTask):
        import pandas as pd
        from ebooklib import epub
        df = pd.read_excel(str(src), sheet_name=0)
        book = epub.EpubBook()
        book.set_identifier("id-1234")
        book.set_title(src.stem)
        book.set_language("es")
        chapter = epub.EpubHtml(title="Contenido", file_name="content.xhtml", lang="es")
        html_content = df.to_html(index=False)
        chapter.content = f"<html><body>{html_content}</body></html>"
        book.add_item(chapter)
        book.spine = ["nav", chapter]
        epub.write_epub(str(dst), book)
        if task.on_progress:
            task.on_progress(1.0)

    def _html_to_csv(self, src: Path, dst: Path, task: ConversionTask):
        import csv
        import re
        from bs4 import BeautifulSoup
        html = src.read_text(encoding="utf-8")
        soup = BeautifulSoup(html, "html.parser")
        rows = []
        for tr in soup.find_all("tr"):
            cells = [cell.get_text(strip=True) for cell in tr.find_all(["td", "th"])]
            if cells:
                rows.append(cells)
        if not rows:
            text = re.sub(r"<[^>]+>", "", html)
            rows = [[line.strip()] for line in text.split("\n") if line.strip()]
        with open(dst, "w", encoding="utf-8", newline="") as f:
            writer = csv.writer(f)
            writer.writerows(rows)
        if task.on_progress:
            task.on_progress(1.0)

    def _html_to_xlsx(self, src: Path, dst: Path, task: ConversionTask):
        import pandas as pd
        from bs4 import BeautifulSoup
        html = src.read_text(encoding="utf-8")
        soup = BeautifulSoup(html, "html.parser")
        rows = []
        for tr in soup.find_all("tr"):
            cells = [cell.get_text(strip=True) for cell in tr.find_all(["td", "th"])]
            if cells:
                rows.append(cells)
        if rows:
            df = pd.DataFrame(rows[1:], columns=rows[0] if len(rows) > 0 else None)
        else:
            df = pd.DataFrame()
        df.to_excel(str(dst), index=False, engine="openpyxl")
        if task.on_progress:
            task.on_progress(1.0)

    def _html_to_odt(self, src: Path, dst: Path, task: ConversionTask):
        import re
        from odf.opendocument import OpenDocumentText
        from odf.text import P
        html = src.read_text(encoding="utf-8")
        text = re.sub(r"<br\s*/?>", "\n", html)
        text = re.sub(r"</p>", "\n", text)
        text = re.sub(r"<[^>]+>", "", text)
        odt = OpenDocumentText()
        for line in text.split("\n"):
            if line.strip():
                odt.text.addElement(P(text=line.strip()))
        odt.save(str(dst))
        if task.on_progress:
            task.on_progress(1.0)

    def _html_to_epub(self, src: Path, dst: Path, task: ConversionTask):
        from ebooklib import epub
        html = src.read_text(encoding="utf-8")
        book = epub.EpubBook()
        book.set_identifier("id-1234")
        book.set_title(src.stem)
        book.set_language("es")
        chapter = epub.EpubHtml(title="Contenido", file_name="content.xhtml", lang="es")
        chapter.content = html
        book.add_item(chapter)
        book.spine = ["nav", chapter]
        epub.write_epub(str(dst), book)
        if task.on_progress:
            task.on_progress(1.0)

    def _rtf_to_csv(self, src: Path, dst: Path, task: ConversionTask):
        import csv
        import re
        rtf = src.read_text(encoding="utf-8", errors="ignore")
        text = re.sub(r"\\par\b", "\n", rtf)
        text = re.sub(r"\{[^}]*\}", "", text)
        text = re.sub(r"\\\w+\d*\s?", "", text)
        text = re.sub(r"[{}]", "", text)
        with open(dst, "w", encoding="utf-8", newline="") as f:
            writer = csv.writer(f)
            for line in text.split("\n"):
                if line.strip():
                    writer.writerow([line.strip()])
        if task.on_progress:
            task.on_progress(1.0)

    def _rtf_to_xlsx(self, src: Path, dst: Path, task: ConversionTask):
        import pandas as pd
        import re
        rtf = src.read_text(encoding="utf-8", errors="ignore")
        text = re.sub(r"\\par\b", "\n", rtf)
        text = re.sub(r"\{[^}]*\}", "", text)
        text = re.sub(r"\\\w+\d*\s?", "", text)
        text = re.sub(r"[{}]", "", text)
        lines = [line.strip() for line in text.split("\n") if line.strip()]
        df = pd.DataFrame({"Contenido": lines})
        df.to_excel(str(dst), index=False, engine="openpyxl")
        if task.on_progress:
            task.on_progress(1.0)

    def _rtf_to_odt(self, src: Path, dst: Path, task: ConversionTask):
        import re
        from odf.opendocument import OpenDocumentText
        from odf.text import P
        rtf = src.read_text(encoding="utf-8", errors="ignore")
        text = re.sub(r"\\par\b", "\n", rtf)
        text = re.sub(r"\{[^}]*\}", "", text)
        text = re.sub(r"\\\w+\d*\s?", "", text)
        text = re.sub(r"[{}]", "", text)
        odt = OpenDocumentText()
        for line in text.split("\n"):
            if line.strip():
                odt.text.addElement(P(text=line.strip()))
        odt.save(str(dst))
        if task.on_progress:
            task.on_progress(1.0)

    def _rtf_to_epub(self, src: Path, dst: Path, task: ConversionTask):
        import re
        from ebooklib import epub
        rtf = src.read_text(encoding="utf-8", errors="ignore")
        text = re.sub(r"\\par\b", "\n", rtf)
        text = re.sub(r"\{[^}]*\}", "", text)
        text = re.sub(r"\\\w+\d*\s?", "", text)
        text = re.sub(r"[{}]", "", text)
        book = epub.EpubBook()
        book.set_identifier("id-1234")
        book.set_title(src.stem)
        book.set_language("es")
        chapter = epub.EpubHtml(title="Contenido", file_name="content.xhtml", lang="es")
        html_content = "<html><body>"
        for line in text.split("\n"):
            if line.strip():
                html_content += f"<p>{line.strip()}</p>"
        html_content += "</body></html>"
        chapter.content = html_content
        book.add_item(chapter)
        book.spine = ["nav", chapter]
        epub.write_epub(str(dst), book)
        if task.on_progress:
            task.on_progress(1.0)

    def _odt_to_csv(self, src: Path, dst: Path, task: ConversionTask):
        import csv
        import zipfile
        from xml.etree import ElementTree as ET
        with zipfile.ZipFile(str(src), "r") as z:
            content_xml = z.read("content.xml")
        root = ET.fromstring(content_xml)
        text_parts = []
        for elem in root.iter():
            if elem.text and elem.text.strip():
                text_parts.append(elem.text.strip())
        with open(dst, "w", encoding="utf-8", newline="") as f:
            writer = csv.writer(f)
            for text in text_parts:
                writer.writerow([text])
        if task.on_progress:
            task.on_progress(1.0)

    def _odt_to_xlsx(self, src: Path, dst: Path, task: ConversionTask):
        import pandas as pd
        import zipfile
        from xml.etree import ElementTree as ET
        with zipfile.ZipFile(str(src), "r") as z:
            content_xml = z.read("content.xml")
        root = ET.fromstring(content_xml)
        text_parts = []
        for elem in root.iter():
            if elem.text and elem.text.strip():
                text_parts.append(elem.text.strip())
        df = pd.DataFrame({"Contenido": text_parts})
        df.to_excel(str(dst), index=False, engine="openpyxl")
        if task.on_progress:
            task.on_progress(1.0)

    def _odt_to_html(self, src: Path, dst: Path, task: ConversionTask):
        import zipfile
        from xml.etree import ElementTree as ET
        with zipfile.ZipFile(str(src), "r") as z:
            content_xml = z.read("content.xml")
        root = ET.fromstring(content_xml)
        text_parts = []
        for elem in root.iter():
            if elem.text and elem.text.strip():
                text_parts.append(elem.text.strip())
        html = f"<!DOCTYPE html><html><head><meta charset='utf-8'><title>{src.stem}</title></head><body>"
        for text in text_parts:
            html += f"<p>{text}</p>"
        html += "</body></html>"
        dst.write_text(html, encoding="utf-8")
        if task.on_progress:
            task.on_progress(1.0)

    def _odt_to_epub(self, src: Path, dst: Path, task: ConversionTask):
        import zipfile
        from xml.etree import ElementTree as ET
        from ebooklib import epub
        with zipfile.ZipFile(str(src), "r") as z:
            content_xml = z.read("content.xml")
        root = ET.fromstring(content_xml)
        text_parts = []
        for elem in root.iter():
            if elem.text and elem.text.strip():
                text_parts.append(elem.text.strip())
        book = epub.EpubBook()
        book.set_identifier("id-1234")
        book.set_title(src.stem)
        book.set_language("es")
        chapter = epub.EpubHtml(title="Contenido", file_name="content.xhtml", lang="es")
        html_content = "<html><body>"
        for text in text_parts:
            html_content += f"<p>{text}</p>"
        html_content += "</body></html>"
        chapter.content = html_content
        book.add_item(chapter)
        book.spine = ["nav", chapter]
        epub.write_epub(str(dst), book)
        if task.on_progress:
            task.on_progress(1.0)

    def _epub_to_csv(self, src: Path, dst: Path, task: ConversionTask):
        import csv
        from ebooklib import epub
        from bs4 import BeautifulSoup
        book = epub.read_epub(str(src))
        items = list(book.get_items_of_type(9))
        rows = []
        for item in items:
            html_content = item.get_content().decode("utf-8")
            soup = BeautifulSoup(html_content, "html.parser")
            for p in soup.find_all(["p", "div"]):
                text = p.get_text(strip=True)
                if text:
                    rows.append([text])
        with open(dst, "w", encoding="utf-8", newline="") as f:
            writer = csv.writer(f)
            writer.writerows(rows)
        if task.on_progress:
            task.on_progress(1.0)

    def _epub_to_xlsx(self, src: Path, dst: Path, task: ConversionTask):
        import pandas as pd
        from ebooklib import epub
        from bs4 import BeautifulSoup
        book = epub.read_epub(str(src))
        items = list(book.get_items_of_type(9))
        rows = []
        for item in items:
            html_content = item.get_content().decode("utf-8")
            soup = BeautifulSoup(html_content, "html.parser")
            for p in soup.find_all(["p", "div"]):
                text = p.get_text(strip=True)
                if text:
                    rows.append(text)
        df = pd.DataFrame({"Contenido": rows})
        df.to_excel(str(dst), index=False, engine="openpyxl")
        if task.on_progress:
            task.on_progress(1.0)

    def _epub_to_odt(self, src: Path, dst: Path, task: ConversionTask):
        from ebooklib import epub
        from bs4 import BeautifulSoup
        from odf.opendocument import OpenDocumentText
        from odf.text import P
        book = epub.read_epub(str(src))
        items = list(book.get_items_of_type(9))
        odt = OpenDocumentText()
        for item in items:
            html_content = item.get_content().decode("utf-8")
            soup = BeautifulSoup(html_content, "html.parser")
            for p in soup.find_all(["p", "div"]):
                text = p.get_text(strip=True)
                if text:
                    odt.text.addElement(P(text=text))
        odt.save(str(dst))
        if task.on_progress:
            task.on_progress(1.0)

    def _epub_to_docx(self, src: Path, dst: Path, task: ConversionTask):
        from ebooklib import epub
        from bs4 import BeautifulSoup
        from docx import Document
        book = epub.read_epub(str(src))
        items = list(book.get_items_of_type(9))
        doc = Document()
        for item in items:
            html_content = item.get_content().decode("utf-8")
            soup = BeautifulSoup(html_content, "html.parser")
            for p in soup.find_all(["p", "div"]):
                text = p.get_text(strip=True)
                if text:
                    doc.add_paragraph(text)
        doc.save(str(dst))
        if task.on_progress:
            task.on_progress(1.0)

    def _pdf_to_html(self, src: Path, dst: Path, task: ConversionTask):
        from pypdf import PdfReader
        reader = PdfReader(str(src))
        html = f"<!DOCTYPE html><html><head><meta charset='utf-8'><title>{src.stem}</title></head><body>"
        for i, page in enumerate(reader.pages):
            text = page.extract_text() or ""
            html += f"<h2>Página {i+1}</h2>"
            for line in text.split("\n"):
                if line.strip():
                    html += f"<p>{line}</p>"
        html += "</body></html>"
        dst.write_text(html, encoding="utf-8")
        if task.on_progress:
            task.on_progress(1.0)

    def _pdf_to_odt(self, src: Path, dst: Path, task: ConversionTask):
        from pypdf import PdfReader
        from odf.opendocument import OpenDocumentText
        from odf.text import P
        reader = PdfReader(str(src))
        odt = OpenDocumentText()
        for page in reader.pages:
            text = page.extract_text() or ""
            for line in text.split("\n"):
                if line.strip():
                    odt.text.addElement(P(text=line))
        odt.save(str(dst))
        if task.on_progress:
            task.on_progress(1.0)

    def _pdf_to_epub(self, src: Path, dst: Path, task: ConversionTask):
        from pypdf import PdfReader
        from ebooklib import epub
        reader = PdfReader(str(src))
        book = epub.EpubBook()
        book.set_identifier("id-1234")
        book.set_title(src.stem)
        book.set_language("es")
        chapter = epub.EpubHtml(title="Contenido", file_name="content.xhtml", lang="es")
        html_content = "<html><body>"
        for i, page in enumerate(reader.pages):
            text = page.extract_text() or ""
            html_content += f"<h2>Página {i+1}</h2>"
            for line in text.split("\n"):
                if line.strip():
                    html_content += f"<p>{line}</p>"
        html_content += "</body></html>"
        chapter.content = html_content
        book.add_item(chapter)
        book.spine = ["nav", chapter]
        epub.write_epub(str(dst), book)
        if task.on_progress:
            task.on_progress(1.0)

    def _pdf_to_xlsx(self, src: Path, dst: Path, task: ConversionTask):
        import pandas as pd
        from pypdf import PdfReader
        reader = PdfReader(str(src))
        rows = []
        for page in reader.pages:
            text = page.extract_text() or ""
            for line in text.split("\n"):
                if line.strip():
                    rows.append([line])
        df = pd.DataFrame(rows, columns=["Contenido"])
        df.to_excel(str(dst), index=False, engine="openpyxl")
        if task.on_progress:
            task.on_progress(1.0)

    def _pptx_to_txt(self, src: Path, dst: Path, task: ConversionTask):
        from pptx import Presentation
        prs = Presentation(str(src))
        text_parts = []
        for i, slide in enumerate(prs.slides):
            text_parts.append(f"--- Slide {i+1} ---")
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        if para.text.strip():
                            text_parts.append(para.text)
        dst.write_text("\n".join(text_parts), encoding="utf-8")
        if task.on_progress:
            task.on_progress(1.0)

    def _pptx_to_docx(self, src: Path, dst: Path, task: ConversionTask):
        from pptx import Presentation
        from docx import Document
        prs = Presentation(str(src))
        doc = Document()
        for i, slide in enumerate(prs.slides):
            doc.add_heading(f"Slide {i+1}", level=1)
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        if para.text.strip():
                            doc.add_paragraph(para.text)
        doc.save(str(dst))
        if task.on_progress:
            task.on_progress(1.0)

    def _pptx_to_html(self, src: Path, dst: Path, task: ConversionTask):
        from pptx import Presentation
        prs = Presentation(str(src))
        html = f"<!DOCTYPE html><html><head><meta charset='utf-8'><title>{src.stem}</title></head><body>"
        for i, slide in enumerate(prs.slides):
            html += f"<h2>Slide {i+1}</h2>"
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        if para.text.strip():
                            html += f"<p>{para.text}</p>"
        html += "</body></html>"
        dst.write_text(html, encoding="utf-8")
        if task.on_progress:
            task.on_progress(1.0)

    def _pptx_to_odt(self, src: Path, dst: Path, task: ConversionTask):
        from pptx import Presentation
        from odf.opendocument import OpenDocumentText
        from odf.text import P, H
        prs = Presentation(str(src))
        odt = OpenDocumentText()
        for i, slide in enumerate(prs.slides):
            odt.text.addElement(H(outlinelevel=1, text=f"Slide {i+1}"))
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        if para.text.strip():
                            odt.text.addElement(P(text=para.text))
        odt.save(str(dst))
        if task.on_progress:
            task.on_progress(1.0)

    def _pptx_to_epub(self, src: Path, dst: Path, task: ConversionTask):
        from pptx import Presentation
        from ebooklib import epub
        prs = Presentation(str(src))
        book = epub.EpubBook()
        book.set_identifier("id-1234")
        book.set_title(src.stem)
        book.set_language("es")
        chapter = epub.EpubHtml(title="Contenido", file_name="content.xhtml", lang="es")
        html_content = "<html><body>"
        for i, slide in enumerate(prs.slides):
            html_content += f"<h2>Slide {i+1}</h2>"
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        if para.text.strip():
                            html_content += f"<p>{para.text}</p>"
        html_content += "</body></html>"
        chapter.content = html_content
        book.add_item(chapter)
        book.spine = ["nav", chapter]
        epub.write_epub(str(dst), book)
        if task.on_progress:
            task.on_progress(1.0)

    def _docx_to_epub(self, src: Path, dst: Path, task: ConversionTask):
        from docx import Document
        from ebooklib import epub
        doc = Document(str(src))
        book = epub.EpubBook()
        book.set_identifier("id-1234")
        book.set_title(src.stem)
        book.set_language("es")
        chapter = epub.EpubHtml(title="Contenido", file_name="content.xhtml", lang="es")
        html_content = "<html><body>"
        for para in doc.paragraphs:
            if para.text.strip():
                html_content += f"<p>{para.text}</p>"
        html_content += "</body></html>"
        chapter.content = html_content
        book.add_item(chapter)
        book.spine = ["nav", chapter]
        epub.write_epub(str(dst), book)
        if task.on_progress:
            task.on_progress(1.0)
