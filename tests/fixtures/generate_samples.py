"""Generate sample test files for integration tests."""
import os
from pathlib import Path


def generate_sample_files(output_dir: Path):
    """Generate all sample test files."""
    output_dir.mkdir(parents=True, exist_ok=True)

    # Documents
    docs_dir = output_dir / "documents"
    docs_dir.mkdir(exist_ok=True)
    _generate_docx(docs_dir / "sample.docx")
    _generate_docx_with_tables(docs_dir / "sample_tables.docx")
    _generate_docx_with_images(docs_dir / "sample_images.docx")
    _generate_pdf(docs_dir / "sample.pdf")
    _generate_xlsx(docs_dir / "sample.xlsx")
    _generate_csv(docs_dir / "sample.csv")
    _generate_txt(docs_dir / "sample.txt")
    _generate_html(docs_dir / "sample.html")
    _generate_rtf(docs_dir / "sample.rtf")
    _generate_pptx(docs_dir / "sample.pptx")

    # Images
    imgs_dir = output_dir / "images"
    imgs_dir.mkdir(exist_ok=True)
    _generate_image(imgs_dir / "sample.png", "RGB", (100, 100), "red")
    _generate_image(imgs_dir / "sample_rgba.png", "RGBA", (100, 100), (255, 0, 0, 128))
    _generate_image(imgs_dir / "sample.jpg", "RGB", (100, 100), "blue")
    _generate_image(imgs_dir / "sample_large.jpg", "RGB", (1000, 1000), "green")

    # Audio
    audio_dir = output_dir / "audio"
    audio_dir.mkdir(exist_ok=True)
    _generate_audio(audio_dir / "sample.wav")

    # Video
    video_dir = output_dir / "video"
    video_dir.mkdir(exist_ok=True)
    _generate_video(video_dir / "sample.mp4")

    print(f"Sample files generated in {output_dir}")


def _generate_docx(path: Path):
    from docx import Document
    doc = Document()
    doc.add_heading("Sample Document", level=1)
    doc.add_paragraph("This is a test paragraph with some content.")
    doc.add_paragraph("Second paragraph with more text for testing conversions.")
    doc.add_paragraph("")
    doc.add_paragraph("Third paragraph after empty line.")
    doc.save(str(path))


def _generate_docx_with_tables(path: Path):
    from docx import Document
    doc = Document()
    doc.add_heading("Document with Tables", level=1)
    doc.add_paragraph("This document contains tables.")

    table = doc.add_table(rows=3, cols=3)
    table.style = "Table Grid"
    headers = ["Name", "Age", "City"]
    for i, header in enumerate(headers):
        table.rows[0].cells[i].text = header

    data = [
        ["Alice", "30", "New York"],
        ["Bob", "25", "London"],
    ]
    for row_idx, row_data in enumerate(data, 1):
        for col_idx, cell_data in enumerate(row_data):
            table.rows[row_idx].cells[col_idx].text = cell_data

    doc.add_paragraph("Text after table.")
    doc.save(str(path))


def _generate_docx_with_images(path: Path):
    from docx import Document
    from docx.shared import Inches
    from PIL import Image
    import io

    doc = Document()
    doc.add_heading("Document with Images", level=1)
    doc.add_paragraph("This document contains an embedded image.")

    img = Image.new("RGB", (200, 200), color="blue")
    img_bytes = io.BytesIO()
    img.save(img_bytes, format="PNG")
    img_bytes.seek(0)

    doc.add_picture(img_bytes, width=Inches(2))
    doc.add_paragraph("Text after image.")
    doc.save(str(path))


def _generate_pdf(path: Path):
    from fpdf import FPDF
    pdf = FPDF()
    pdf.add_page()
    pdf.set_font("Helvetica", size=12)
    pdf.cell(0, 10, "Sample PDF Document", ln=True)
    pdf.ln(5)
    pdf.multi_cell(0, 5, "This is a test PDF file with some content for testing conversions.")
    pdf.ln(5)
    pdf.multi_cell(0, 5, "Second paragraph with more text.")
    pdf.output(str(path))


def _generate_xlsx(path: Path):
    import pandas as pd
    data = {
        "Name": ["Alice", "Bob", "Charlie"],
        "Age": [30, 25, 35],
        "City": ["New York", "London", "Paris"],
        "Salary": [50000, 60000, 75000],
    }
    df = pd.DataFrame(data)
    df.to_excel(str(path), index=False, engine="openpyxl")


def _generate_csv(path: Path):
    import csv
    with open(path, "w", newline="", encoding="utf-8") as f:
        writer = csv.writer(f)
        writer.writerow(["Name", "Age", "City", "Salary"])
        writer.writerow(["Alice", "30", "New York", "50000"])
        writer.writerow(["Bob", "25", "London", "60000"])
        writer.writerow(["Charlie", "35", "Paris", "75000"])


def _generate_txt(path: Path):
    content = """Sample Text File

This is a test text file with multiple paragraphs.

Second paragraph with more content for testing conversions.

Third paragraph after empty line.

Final paragraph with some text."""
    path.write_text(content, encoding="utf-8")


def _generate_html(path: Path):
    html = """<!DOCTYPE html>
<html>
<head><meta charset="utf-8"><title>Sample HTML</title></head>
<body>
<h1>Sample HTML Document</h1>
<p>This is a test HTML file with some content.</p>
<p>Second paragraph with more text for testing.</p>
<ul>
<li>Item 1</li>
<li>Item 2</li>
<li>Item 3</li>
</ul>
</body>
</html>"""
    path.write_text(html, encoding="utf-8")


def _generate_rtf(path: Path):
    rtf = """{\\rtf1\\ansi\\deff0
{\\fonttbl{\\f0 Times New Roman;}}
\\pard\\qc\\b\\fs32 Sample RTF Document\\b0\\par
\\pard\\fs24 This is a test RTF file with some content.\\par
\\par
Second paragraph with more text for testing.\\par
\\par
{\\b Bold text} and {\\i italic text} in RTF format.\\par
}"""
    path.write_text(rtf, encoding="utf-8")


def _generate_pptx(path: Path):
    from pptx import Presentation
    prs = Presentation()

    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "Sample Presentation"
    slide.placeholders[1].text = "Test content for PPTX conversion"

    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    slide2.shapes.title.text = "Slide 2"
    slide2.placeholders[1].text = "More test content"

    prs.save(str(path))


def _generate_image(path: Path, mode: str, size: tuple, color):
    from PIL import Image
    img = Image.new(mode, size, color=color)
    img.save(str(path))


def _generate_audio(path: Path):
    import wave
    import struct
    with wave.open(str(path), "w") as wf:
        wf.setnchannels(1)
        wf.setsampwidth(2)
        wf.setframerate(44100)
        for i in range(44100):
            value = int(32767 * 0.5 * (i % 100 < 50))
            wf.writeframes(struct.pack("<h", value))


def _generate_video(path: Path):
    try:
        import subprocess
        cmd = [
            "ffmpeg", "-y", "-f", "lavfi", "-i",
            "color=c=blue:s=320x240:d=1",
            "-f", "lavfi", "-i", "anullsrc=r=44100:cl=mono",
            "-t", "1", "-shortest", str(path)
        ]
        subprocess.run(cmd, capture_output=True, timeout=10)
    except Exception:
        print(f"Warning: Could not generate sample video (FFmpeg not available)")


if __name__ == "__main__":
    generate_sample_files(Path(__file__).parent / "sample_files")
